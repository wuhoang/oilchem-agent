"""
Office 文档读写工具。

支持 Excel (.xlsx)、Word (.docx)、PowerPoint (.pptx) 的读取和写入，
让 Agent 能够查看、分析、生成各种 Office 文件。
"""

from __future__ import annotations

import os
from typing import Any

from loguru import logger

from app.tools.base import BaseTool, ToolMetadata, ToolResult
from app.tools.registry import register_tool

# 工作目录
WORK_DIR = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", "..", "..", "..", ".."))


def _resolve_path(path: str) -> str:
    """将相对路径解析为绝对路径。"""
    if os.path.isabs(path):
        return path
    return os.path.join(WORK_DIR, path)


# ---------------------------------------------------------------------------
# Excel 工具
# ---------------------------------------------------------------------------

@register_tool(ToolMetadata(
    name="read_excel",
    category="office",
    description=(
        "读取 Excel 文件（.xlsx）的内容，返回工作表数据。"
        "支持指定 sheet_name 读取特定工作表，或读取所有工作表。"
        "返回格式为 JSON 数组，每个元素是一行数据（字典）。"
    ),
    parameters={
        "file_path": {
            "type": "string",
            "description": "Excel 文件路径，如 'data/report.xlsx' 或绝对路径",
        },
        "sheet_name": {
            "type": "string",
            "description": "工作表名称（可选），不指定则读取第一个工作表",
        },
        "sheet_index": {
            "type": "integer",
            "description": "工作表索引（可选），从 0 开始，如 0 表示第一个工作表",
        },
        "header_row": {
            "type": "integer",
            "description": "表头所在行号（0开始），默认为 0",
        },
        "max_rows": {
            "type": "integer",
            "description": "最大读取行数（可选），用于限制大数据集",
        },
    },
))
class ReadExcelTool(BaseTool):
    """读取 Excel 文件。"""

    async def execute(self, **kwargs: Any) -> ToolResult:
        file_path = kwargs.get("file_path", "").strip()
        if not file_path:
            return ToolResult(success=False, error="缺少 file_path 参数")

        abs_path = _resolve_path(file_path)
        if not os.path.exists(abs_path):
            return ToolResult(success=False, error=f"文件不存在: {abs_path}")

        try:
            import openpyxl

            wb = openpyxl.load_workbook(abs_path, data_only=True)
            sheet_name = kwargs.get("sheet_name")
            sheet_index = kwargs.get("sheet_index")

            if sheet_name:
                if sheet_name not in wb.sheetnames:
                    return ToolResult(
                        success=False,
                        error=f"工作表 '{sheet_name}' 不存在，可用: {wb.sheetnames}",
                    )
                ws = wb[sheet_name]
            elif sheet_index is not None:
                if sheet_index >= len(wb.sheetnames):
                    return ToolResult(
                        success=False,
                        error=f"工作表索引 {sheet_index} 超出范围，共 {len(wb.sheetnames)} 个",
                    )
                ws = wb[wb.sheetnames[sheet_index]]
            else:
                ws = wb.active

            header_row = kwargs.get("header_row", 0)
            max_rows = kwargs.get("max_rows")

            rows = list(ws.iter_rows(min_row=1, values_only=True))
            if not rows:
                return ToolResult(success=True, data={
                    "file": file_path,
                    "sheet": ws.title,
                    "headers": [],
                    "data": [],
                    "row_count": 0,
                })

            # 提取表头
            headers = [str(h) if h is not None else f"col_{i}" for i, h in enumerate(rows[header_row])]

            # 提取数据行
            data = []
            for row_idx, row in enumerate(rows[header_row + 1:], start=header_row + 2):
                record = {}
                for col_idx, val in enumerate(row):
                    if col_idx < len(headers):
                        record[headers[col_idx]] = val
                data.append(record)
                if max_rows and len(data) >= max_rows:
                    break

            return ToolResult(success=True, data={
                "file": file_path,
                "sheet": ws.title,
                "sheets": wb.sheetnames,
                "headers": headers,
                "data": data,
                "row_count": len(data),
                "total_rows": ws.max_row,
                "total_columns": ws.max_column,
            })

        except Exception as exc:
            logger.bind(component="office").error("read_excel failed: {}", exc)
            return ToolResult(success=False, error=f"读取 Excel 失败: {exc}")


@register_tool(ToolMetadata(
    name="write_excel",
    category="office",
    description=(
        "将数据写入 Excel 文件。支持创建新文件或追加到现有文件。"
        "参数 data 为数组，每个元素是一行数据的字典。"
        "headers 指定列名（可选）。"
    ),
    parameters={
        "file_path": {
            "type": "string",
            "description": "输出文件路径，如 'output/report.xlsx'",
        },
        "data": {
            "type": "array",
            "description": "要写入的数据数组，每个元素是对象，如 [{'name': '张三', 'score': 95}]",
        },
        "sheet_name": {
            "type": "string",
            "description": "工作表名称（可选），默认 'Sheet1'",
        },
        "headers": {
            "type": "array",
            "description": "列名数组（可选），不指定则从数据对象的键提取",
        },
        "append": {
            "type": "boolean",
            "description": "是否追加到现有文件（默认 false，覆盖写入）",
        },
    },
))
class WriteExcelTool(BaseTool):
    """写入 Excel 文件。"""

    async def execute(self, **kwargs: Any) -> ToolResult:
        file_path = kwargs.get("file_path", "").strip()
        data = kwargs.get("data", [])

        if not file_path:
            return ToolResult(success=False, error="缺少 file_path 参数")
        if not data:
            return ToolResult(success=False, error="缺少 data 参数")

        abs_path = _resolve_path(file_path)
        sheet_name = kwargs.get("sheet_name", "Sheet1")
        headers = kwargs.get("headers")
        append = kwargs.get("append", False)

        try:
            import openpyxl

            if append and os.path.exists(abs_path):
                wb = openpyxl.load_workbook(abs_path)
                if sheet_name in wb.sheetnames:
                    ws = wb[sheet_name]
                else:
                    ws = wb.create_sheet(sheet_name)
            else:
                wb = openpyxl.Workbook()
                ws = wb.active
                ws.title = sheet_name

            # 确定列名
            if not headers:
                if data:
                    headers = list(data[0].keys())
                else:
                    headers = []

            # 检查是否需要写表头
            start_row = 1
            if ws.max_row == 1 or not append:
                for col_idx, header in enumerate(headers, start=1):
                    ws.cell(row=1, column=col_idx, value=header)
                start_row = 2

            # 写入数据
            for row_idx, record in enumerate(data, start=start_row):
                for col_idx, header in enumerate(headers, start=1):
                    ws.cell(
                        row=row_idx,
                        column=col_idx,
                        value=record.get(header, "")
                    )

            # 确保目录存在
            os.makedirs(os.path.dirname(abs_path) or ".", exist_ok=True)
            wb.save(abs_path)

            return ToolResult(success=True, data={
                "file": file_path,
                "sheet": sheet_name,
                "rows_written": len(data),
                "columns": len(headers),
                "message": f"成功写入 {len(data)} 行数据到 {file_path}",
            })

        except Exception as exc:
            logger.bind(component="office").error("write_excel failed: {}", exc)
            return ToolResult(success=False, error=f"写入 Excel 失败: {exc}")


# ---------------------------------------------------------------------------
# Word 工具
# ---------------------------------------------------------------------------

@register_tool(ToolMetadata(
    name="read_word",
    category="office",
    description=(
        "读取 Word 文档（.docx）的内容，返回段落列表和表格数据。"
        "可用于提取文档文本、分析报告内容等。"
    ),
    parameters={
        "file_path": {
            "type": "string",
            "description": "Word 文件路径，如 'data/report.docx'",
        },
        "include_tables": {
            "type": "boolean",
            "description": "是否包含表格数据（默认 true）",
        },
        "max_paragraphs": {
            "type": "integer",
            "description": "最大读取段落数（可选）",
        },
    },
))
class ReadWordTool(BaseTool):
    """读取 Word 文档。"""

    async def execute(self, **kwargs: Any) -> ToolResult:
        file_path = kwargs.get("file_path", "").strip()
        if not file_path:
            return ToolResult(success=False, error="缺少 file_path 参数")

        abs_path = _resolve_path(file_path)
        if not os.path.exists(abs_path):
            return ToolResult(success=False, error=f"文件不存在: {abs_path}")

        try:
            from docx import Document

            doc = Document(abs_path)
            include_tables = kwargs.get("include_tables", True)
            max_paragraphs = kwargs.get("max_paragraphs")

            paragraphs = []
            for i, para in enumerate(doc.paragraphs):
                text = para.text.strip()
                if text:
                    paragraphs.append({
                        "index": i,
                        "style": para.style.name if para.style else "Normal",
                        "text": text,
                    })
                if max_paragraphs and len(paragraphs) >= max_paragraphs:
                    break

            # 提取标题
            headings = [p for p in paragraphs if p["style"].startswith("Heading")]

            # 提取表格
            tables = []
            if include_tables:
                for t_idx, table in enumerate(doc.tables):
                    rows_data = []
                    for row in table.rows:
                        row_data = [cell.text.strip() for cell in row.cells]
                        rows_data.append(row_data)
                    tables.append({
                        "index": t_idx,
                        "rows": len(table.rows),
                        "columns": len(table.columns),
                        "data": rows_data,
                    })

            # 合并全文
            full_text = "\n".join(p["text"] for p in paragraphs)

            return ToolResult(success=True, data={
                "file": file_path,
                "paragraphs": paragraphs,
                "headings": headings,
                "tables": tables,
                "full_text": full_text,
                "total_paragraphs": len(doc.paragraphs),
                "total_tables": len(doc.tables),
            })

        except Exception as exc:
            logger.bind(component="office").error("read_word failed: {}", exc)
            return ToolResult(success=False, error=f"读取 Word 文档失败: {exc}")


@register_tool(ToolMetadata(
    name="write_word",
    category="office",
    description=(
        "创建 Word 文档（.docx）。支持写入标题、段落和表格。"
        "参数 paragraphs 为数组，每个元素包含 text 和 style。"
        "参数 tables 为二维数组，每个子数组是一张表。"
    ),
    parameters={
        "file_path": {
            "type": "string",
            "description": "输出文件路径，如 'output/report.docx'",
        },
        "title": {
            "type": "string",
            "description": "文档标题（可选）",
        },
        "paragraphs": {
            "type": "array",
            "description": "段落数组，每个元素: {'text': '内容', 'style': 'Normal'} 或 {'text': '标题', 'style': 'Heading 1'}",
        },
        "tables": {
            "type": "array",
            "description": "表格二维数组，如 [[['姓名','分数'],['张三',95],['李四',87]]]",
        },
    },
))
class WriteWordTool(BaseTool):
    """创建 Word 文档。"""

    async def execute(self, **kwargs: Any) -> ToolResult:
        file_path = kwargs.get("file_path", "").strip()
        if not file_path:
            return ToolResult(success=False, error="缺少 file_path 参数")

        abs_path = _resolve_path(file_path)
        title = kwargs.get("title", "")
        paragraphs = kwargs.get("paragraphs", [])
        tables = kwargs.get("tables", [])

        try:
            from docx import Document

            doc = Document()

            # 标题
            if title:
                doc.add_heading(title, level=0)

            # 段落
            for para in paragraphs:
                if isinstance(para, dict):
                    text = para.get("text", "")
                    style = para.get("style", "Normal")
                    if style.startswith("Heading"):
                        level = int(style.split()[-1]) if style.split()[-1].isdigit() else 1
                        doc.add_heading(text, level=min(level, 9))
                    else:
                        doc.add_paragraph(text, style=style)
                elif isinstance(para, str):
                    doc.add_paragraph(para)

            # 表格
            for table_data in tables:
                if isinstance(table_data, list) and len(table_data) > 0:
                    rows_count = len(table_data)
                    cols_count = len(table_data[0]) if table_data[0] else 0
                    if rows_count > 0 and cols_count > 0:
                        table = doc.add_table(rows=rows_count, cols=cols_count)
                        table.style = "Table Grid"
                        for i, row_data in enumerate(table_data):
                            for j, cell_value in enumerate(row_data):
                                if j < cols_count:
                                    table.rows[i].cells[j].text = str(cell_value)

            # 确保目录存在
            os.makedirs(os.path.dirname(abs_path) or ".", exist_ok=True)
            doc.save(abs_path)

            return ToolResult(success=True, data={
                "file": file_path,
                "paragraphs_written": len(paragraphs),
                "tables_written": len(tables),
                "message": f"成功创建 Word 文档: {file_path}",
            })

        except Exception as exc:
            logger.bind(component="office").error("write_word failed: {}", exc)
            return ToolResult(success=False, error=f"创建 Word 文档失败: {exc}")


# ---------------------------------------------------------------------------
# PowerPoint 工具
# ---------------------------------------------------------------------------

@register_tool(ToolMetadata(
    name="read_ppt",
    category="office",
    description=(
        "读取 PowerPoint 文件（.pptx）的内容，返回幻灯片文本和备注。"
        "可用于提取演示文稿内容、分析报告结构等。"
    ),
    parameters={
        "file_path": {
            "type": "string",
            "description": "PPT 文件路径，如 'data/presentation.pptx'",
        },
        "slide_index": {
            "type": "integer",
            "description": "读取指定幻灯片索引（可选），从 0 开始",
        },
    },
))
class ReadPPTTool(BaseTool):
    """读取 PowerPoint 文件。"""

    async def execute(self, **kwargs: Any) -> ToolResult:
        file_path = kwargs.get("file_path", "").strip()
        if not file_path:
            return ToolResult(success=False, error="缺少 file_path 参数")

        abs_path = _resolve_path(file_path)
        if not os.path.exists(abs_path):
            return ToolResult(success=False, error=f"文件不存在: {abs_path}")

        try:
            from pptx import Presentation

            prs = Presentation(abs_path)
            slide_index = kwargs.get("slide_index")

            slides_data = []
            for idx, slide in enumerate(prs.slides):
                if slide_index is not None and idx != slide_index:
                    continue

                shapes_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        shapes_text.append(shape.text.strip())

                # 提取表格数据
                tables_data = []
                for shape in slide.shapes:
                    if shape.has_table:
                        table = shape.table
                        table_data = []
                        for row in table.rows:
                            row_data = [cell.text for cell in row.cells]
                            table_data.append(row_data)
                        tables_data.append(table_data)

                notes_text = ""
                if slide.has_notes_slide:
                    notes_text = slide.notes_slide.notes_text_frame.text.strip()

                slides_data.append({
                    "index": idx,
                    "text_content": shapes_text,
                    "tables": tables_data,
                    "notes": notes_text,
                })

            all_text = []
            for s in slides_data:
                all_text.extend(s["text_content"])

            return ToolResult(success=True, data={
                "file": file_path,
                "total_slides": len(prs.slides),
                "slides": slides_data,
                "full_text": "\n".join(all_text),
            })

        except Exception as exc:
            logger.bind(component="office").error("read_ppt failed: {}", exc)
            return ToolResult(success=False, error=f"读取 PPT 失败: {exc}")


@register_tool(ToolMetadata(
    name="write_ppt",
    category="office",
    description=(
        "创建 PowerPoint 文件（.pptx）。支持添加幻灯片、设置标题和内容。"
        "参数 slides 为数组，每个元素包含 title、content 和 notes。"
    ),
    parameters={
        "file_path": {
            "type": "string",
            "description": "输出文件路径，如 'output/presentation.pptx'",
        },
        "slides": {
            "type": "array",
            "description": "幻灯片数组，每个元素: {'title': '标题', 'content': ['要点1', '要点2'], 'notes': '备注'}",
        },
        "template": {
            "type": "string",
            "description": "模板类型（可选）: 'blank'(空白), 'title'(标题页), 'content'(标题+内容)，默认 'content'",
        },
    },
))
class WritePPTTool(BaseTool):
    """创建 PowerPoint 文件。"""

    async def execute(self, **kwargs: Any) -> ToolResult:
        file_path = kwargs.get("file_path", "").strip()
        if not file_path:
            return ToolResult(success=False, error="缺少 file_path 参数")

        abs_path = _resolve_path(file_path)
        slides = kwargs.get("slides", [])
        template = kwargs.get("template", "content")

        if not slides:
            return ToolResult(success=False, error="缺少 slides 参数")

        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt

            prs = Presentation()

            # 选择布局
            if template == "title":
                slide_layout = prs.slide_layouts[0]  # Title Slide
            elif template == "blank":
                slide_layout = prs.slide_layouts[6]  # Blank
            else:
                slide_layout = prs.slide_layouts[1]  # Title and Content

            for slide_data in slides:
                if isinstance(slide_data, dict):
                    title = slide_data.get("title", "")
                    content = slide_data.get("content", [])
                    notes = slide_data.get("notes", "")

                    slide = prs.slides.add_slide(slide_layout)

                    # 设置标题
                    if title and slide.shapes.title:
                        slide.shapes.title.text = title

                    # 设置内容
                    if content:
                        for shape in slide.shapes:
                            if shape.has_text_frame and shape != slide.shapes.title:
                                tf = shape.text_frame
                                tf.clear()
                                for i, item in enumerate(content):
                                    if i == 0:
                                        tf.paragraphs[0].text = str(item)
                                    else:
                                        p = tf.add_paragraph()
                                        p.text = str(item)
                                break

                    # 添加备注
                    if notes:
                        notes_slide = slide.notes_slide
                        notes_slide.notes_text_frame.text = notes

            # 确保目录存在
            os.makedirs(os.path.dirname(abs_path) or ".", exist_ok=True)
            prs.save(abs_path)

            return ToolResult(success=True, data={
                "file": file_path,
                "slides_created": len(slides),
                "message": f"成功创建 {len(slides)} 页 PPT: {file_path}",
            })

        except Exception as exc:
            logger.bind(component="office").error("write_ppt failed: {}", exc)
            return ToolResult(success=False, error=f"创建 PPT 失败: {exc}")


__all__ = [
    "ReadExcelTool",
    "WriteExcelTool",
    "ReadWordTool",
    "WriteWordTool",
    "ReadPPTTool",
    "WritePPTTool",
]
